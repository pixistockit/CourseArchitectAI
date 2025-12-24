# /modules/audit_slide/analyzer.py

import re
import os
import csv
import json
import importlib
from datetime import datetime
from spellchecker import SpellChecker
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.dml import MSO_FILL

# --- 1. CRITICAL CONFIG IMPORTS ---
try:
    from . import config as config_module
    importlib.reload(config_module)

    from .config import (
        EXEMPT_SHAPE_NAMES, TITLE_FONT_NAME, TITLE_FONT_SIZE_MIN, TITLE_FONT_SIZE_MAX,
        TITLE_MUST_BE_BOLD, BODY_FONT_SIZE_MIN, ALLOWED_BODY_FONTS,
        NOTES_FONT_NAME, NOTES_FONT_SIZE_MIN, NOTES_FONT_SIZE_MAX, NOTES_FONT_COLOR_RGB,
        THEME_FONT_MAPPING, WCAG_MIN_GRAPHIC_RATIO, REQUIRED_HEADERS, VALID_GAGNE_TERMS, GAGNE_CATEGORIES,
        EXEMPT_FIRST_SLIDE, EXEMPT_LAST_SLIDE, EXEMPT_SPECIFIC_SLIDES, ACTIVITY_KEYWORDS
    )
except ImportError as e:
    print(f"❌ Critical Config Error: {e}")
    REQUIRED_HEADERS = []
    NOTES_FONT_NAME = None

# --- 2. CRITICAL UTILS IMPORTS ---
try:
    from .utils import (
        rgb_pptx_to_tuple, shapes_overlap, get_required_ratio, calculate_contrast_ratio, 
        is_brand_color, find_closest_compliant_color
    )
except ImportError as e:
    print(f"❌ Critical Utils Error: {e}")
    # Safety fallbacks
    def shapes_overlap(a, b): return False
    def rgb_pptx_to_tuple(c): return (0,0,0)
    def calculate_contrast_ratio(a, b): return 21.0
    def is_brand_color(c): return True
    def get_required_ratio(x): return 4.5

# --- 3. OPTIONAL AI MODULES ---
try:
    from .grammar_master import ClarityEngine
    from .time_estimator import TimeEstimator
except ImportError:
    print("⚠️ Analyzer warning: Optional AI modules not found.")
    ClarityEngine = None
    TimeEstimator = None

# Constants
BG_IS_IMAGE = -1
BG_IS_TABLE = -2
WHITE_RGB_TUPLE = (255, 255, 255)
BLACK_RGB_TUPLE = (0, 0, 0)

# --- HELPER: RGB TO HEX ---
def rgb_to_hex(rgb_tuple):
    """Converts (255, 255, 255) -> #FFFFFF"""
    if not rgb_tuple or len(rgb_tuple) != 3: return "???"
    return "#{:02x}{:02x}{:02x}".format(*rgb_tuple).upper()

# --- HYBRID COMPATIBILITY CLASS ---
class HybridReport(dict):
    """
    Acts as a Dictionary for the new dashboard (accessing ['summary']),
    but acts as a List for legacy tools (iterating over issues).
    """
    def __iter__(self):
        return iter(self.get('detailed_issues', []))

class PptxAnalyzer:
    def __init__(self, pptx_path: str):
        self.pptx_path = pptx_path
        self.filename = os.path.basename(pptx_path)
        self.prs = Presentation(pptx_path)
        self.total_slides_checked = 0
        self.qa_report = []
        
        # --- DYNAMIC SETTINGS LOAD ---
        self.active_buffer = 0.0
        json_config_path = os.path.join('data', 'config', 'llm_config.json')
        
        config_loaded = False
        if os.path.exists(json_config_path):
            try:
                with open(json_config_path, 'r') as f:
                    user_config = json.load(f)
                    if 'default_buffer' in user_config and str(user_config['default_buffer']).strip():
                        self.active_buffer = float(user_config['default_buffer'])
                        config_loaded = True
            except Exception as e:
                print(f"⚠️ JSON Config Load Error: {e}")

        if not config_loaded:
            try:
                if hasattr(config_module, 'default_buffer'):
                    self.active_buffer = float(config_module.default_buffer)
                elif hasattr(config_module, 'BUFFER_ACTIVITY_SLIDE'):
                    self.active_buffer = float(config_module.BUFFER_ACTIVITY_SLIDE)
            except:
                self.active_buffer = 0.0

        self.gagne_metrics = { 
            "Gain Attention": 0, "Inform Objectives": 0, "Stimulate Recall": 0,
            "Present Content": 0, "Provide Guidance": 0, "Elicit Performance": 0,
            "Provide Feedback": 0, "Assess Performance": 0, "Enhance Retention": 0,
            "Other": 0 
        }
        
        self.gagne_time_distribution = {k: 0.0 for k in self.gagne_metrics}

        self.pacing_data = {
            "total_planned_mins": 0, 
            "total_estimated_mins": 0, 
            "total_projected_mins": 0, 
            "used_buffer": self.active_buffer, 
            "sections": [] 
        }
        self.current_section = {"name": "Introduction", "planned": 0, "estimated": 0, "projected": 0, "slide_count": 0}
        
        self.debug_log = [] 
        self.slide_content_map = [] 

        self.spell = SpellChecker()
        self.spell.word_frequency.load_words(['pptx', 'wcag', 'gagne', 'rgb', 'id', 'qa', 'calibri', 'rockwell', "today's", "it's", "don't", "can't", "you'll"])
        
        self.clarity = ClarityEngine() if ClarityEngine else None
        self.timer = TimeEstimator() if TimeEstimator else None

    # --- MAP GAGNE EVENTS ---
    def _map_gagne_event(self, text, slide_title=""):
        t = text.lower()
        title_lower = slide_title.lower()
        events = []

        if re.search(r'knowledge\s*check|quiz|test|assessment|exam', title_lower):
            events.append("Assess Performance")

        target_line = ""
        found_strict_header = False
        lines = text.split('\n')
        for line in lines:
            if re.match(r'(?:instructional\s+)?(?:activity|gagne\s*event)\s*:', line, re.IGNORECASE):
                target_line = line.lower()
                found_strict_header = True
                break
        
        if not found_strict_header:
            for line in lines:
                if len(line) < 100 and len(line.strip()) > 3:
                    if self._check_event_keywords(line):
                        target_line = line.lower()
                        break
        
        analysis_text = target_line if target_line else ""
        
        if analysis_text:
            if re.search(r'gain\s*attention|hook|ice\s*breaker|warm\s*up|welcome', analysis_text): events.append("Gain Attention")
            if re.search(r'inform\s*objectives|objective|outcome|goal|agenda|direction|overview|roadmap', analysis_text): events.append("Inform Objectives")
            if re.search(r'stimulate\s*recall|recall|review|prior\s*knowledge|refresh|remind', analysis_text): events.append("Stimulate Recall")
            if re.search(r'present\s*content|presentation|lecture|explain|demonstrate|show|teach|concept', analysis_text): events.append("Present Content")
            if re.search(r'provide\s*guidance|guidance|guided\s*learning|scaffold|support|coach|help|tip', analysis_text): events.append("Provide Guidance")
            if re.search(r'elicit\s*performance|practice|group\s*activity|learner\s*activity|exercise|simulation|role\s*play|hands\s*on|worksheet', analysis_text): events.append("Elicit Performance")
            if re.search(r'provide\s*feedback|feedback|debrief|review\s*answer|correct|discussion', analysis_text): events.append("Provide Feedback")
            if re.search(r'assess\s*performance|assessment|quiz|test|check|knowledge\s*check|evaluate', analysis_text): events.append("Assess Performance")
            if re.search(r'enhance\s*retention|retention|transfer|wrap\s*up|summary|conclusion|close|job\s*aid|takeaway', analysis_text): events.append("Enhance Retention")

        return list(set(events)) if events else []

    def _check_event_keywords(self, text):
        t = text.lower()
        return bool(re.search(r'gain\s*attention|inform\s*objectives|stimulate\s*recall|present\s*content|provide\s*guidance|elicit\s*performance|provide\s*feedback|assess\s*performance|enhance\s*retention', t))

    def _calculate_cadence_metrics(self):
        for k in self.gagne_metrics: 
            self.gagne_metrics[k] = 0
            self.gagne_time_distribution[k] = 0.0

        self.pacing_data["total_planned_mins"] = 0
        self.pacing_data["total_projected_mins"] = 0
        self.pacing_data["sections"] = []
        self.current_section = {"name": "Introduction", "planned": 0, "estimated": 0, "projected": 0, "slide_count": 0}

        active_funded_activity = "" 

        for slide in self.slide_content_map:
            notes = (slide['notes'] or "")
            full_text = (slide['full_text'] or "").lower()
            title = (slide['title'] or "")
            layout_name = slide['visual_context']['layout'].lower()

            is_section_break = "section" in layout_name or "agenda" in layout_name or "divider" in layout_name
            if not is_section_break and re.match(r'\d+\.\d+\.\d+', title):
                is_section_break = True

            if is_section_break:
                active_funded_activity = "" 
                if self.current_section["slide_count"] > 0:
                    self.pacing_data["sections"].append(self.current_section.copy())
                section_name = title if title else f"Section {slide['slide_number']}"
                self.current_section = {"name": section_name, "planned": 0, "estimated": 0, "projected": 0, "slide_count": 0}

            explicit_time = 0
            time_match = re.search(r'(?:time|duration):\s*(\d+)', notes.lower())
            if time_match: explicit_time = int(time_match.group(1))

            act_match = re.search(r'activity:\s*([^\n]+)', notes.lower())
            current_activity = act_match.group(1).strip().lower() if act_match else ""

            detected_events = self._map_gagne_event(notes, title)
            is_interactive = any(ev in ["Provide Guidance", "Elicit Performance", "Assess Performance"] for ev in detected_events)

            slide_time = 0
            
            if explicit_time > 0:
                slide_time = explicit_time
                active_funded_activity = current_activity
            else:
                if current_activity and current_activity == active_funded_activity:
                    slide_time = 0
                else:
                    word_count = len((notes + " " + full_text).split())
                    word_count_time = max(0.5, word_count / 130)
                    
                    if is_interactive and current_activity and self.active_buffer > 0:
                        slide_time = max(word_count_time, self.active_buffer)
                    else:
                        slide_time = word_count_time
                    
                    active_funded_activity = "" 

            self.pacing_data["total_planned_mins"] += explicit_time
            self.pacing_data["total_projected_mins"] += slide_time
            
            self.current_section["slide_count"] += 1
            self.current_section["projected"] += slide_time
            self.current_section["planned"] += explicit_time

            if detected_events:
                total_weight = 0
                weights = {}
                for ev in detected_events:
                    w = 1
                    if ev in ["Elicit Performance", "Assess Performance"]: w = 10
                    elif ev in ["Present Content", "Provide Feedback"]: w = 3
                    weights[ev] = w
                    total_weight += w
                
                for ev in detected_events:
                    share = (weights[ev] / total_weight) * slide_time
                    self.gagne_metrics[ev] += 1
                    self.gagne_time_distribution[ev] += share
            else:
                self.gagne_metrics["Other"] += 1

        if self.current_section["slide_count"] > 0:
            self.pacing_data["sections"].append(self.current_section)

    def export_cadence_log(self, output_dir):
        pass

    # ... (Helpers) ...
    def _extract_slide_text_content(self, slide):
        title = ""
        text_parts = []
        if slide.shapes.title and slide.shapes.title.text:
            title = slide.shapes.title.text.strip()
        img_count = 0; chart_count = 0; table_count = 0; group_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                clean_t = shape.text.strip()
                if clean_t and clean_t != title: text_parts.append(clean_t)
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE: img_count += 1
            if shape.has_chart: chart_count += 1
            if shape.has_table: table_count += 1
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP: group_count += 1
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        layout_name = slide.slide_layout.name if slide.slide_layout else "Unknown"
        return {
            "title": title, "text": " ".join(text_parts), "notes": notes,
            "visual_context": { "layout": layout_name, "images": img_count, "charts": chart_count, "tables": table_count, "groups": group_count }
        }

    def _strip_citations_for_analysis(self, text):
        if not text: return ""
        text_lower = text.lower()
        headers = [r"source\s*:", r"sources\s*:", r"references?\s*:", r"bibliography\s*:", r"works\s+cited\s*:"]
        best_idx = len(text); found = False
        for h in headers:
            match = re.search(h, text_lower)
            if match:
                for m in re.finditer(h, text_lower):
                    if m.start() > len(text) * 0.5: 
                        best_idx = m.start(); found = True
        if found: return text[:best_idx].strip()
        split_point = int(len(text) * 0.8)
        bottom_chunk = text_lower[split_point:]
        if "http" in bottom_chunk:
            idx = bottom_chunk.find("http")
            return text[:split_point + idx].strip()
        return text

    def _check_single_quotes(self, text, shape_name, slide_index, slide_issues):
        if not text: return
        double_quote_pattern = r'([“"])(.*?)\1'
        safe_zones = []
        for match in re.finditer(double_quote_pattern, text, re.DOTALL): safe_zones.append(match.span())
        single_quote_pattern = r"(?<!\w)['‘](.+?)['’](?!\w)"
        offenders = []
        for match in re.finditer(single_quote_pattern, text, re.DOTALL):
            s_start, s_end = match.span(); phrase = match.group(0); is_safe = False
            for (d_start, d_end) in safe_zones:
                if d_start <= s_start and s_end <= d_end: is_safe = True; break
            if not is_safe:
                clean_phrase = phrase.strip()
                if len(clean_phrase) > 2: offenders.append(clean_phrase)
        if offenders:
            details_str = ", ".join(offenders)
            slide_issues.append({"slide": slide_index, "check": "Style - Punctuation", "shape_name": shape_name, "result": "FAIL", "details": f"Do not use single quotes for emphasis. Found: {details_str}"})

    def _resolve_font_name(self, font_name: str) -> str:
        if not font_name: return "Calibri"
        normalized_name = font_name.lower().strip()
        return THEME_FONT_MAPPING.get(normalized_name, font_name)

    def _is_exempt_shape(self, shape) -> bool:
        s_name = shape.name.lower()
        if shape.is_placeholder:
            ph_type = shape.placeholder_format.type
            if ph_type in (PP_PLACEHOLDER.SLIDE_NUMBER, PP_PLACEHOLDER.FOOTER, PP_PLACEHOLDER.DATE): return True
        if any(name in s_name for name in EXEMPT_SHAPE_NAMES): return True
        if shape.has_text_frame and shape.text_frame.text and shape.text_frame.text.strip():
            text_lower = shape.text_frame.text.lower()
            if "copyright" in text_lower or "rights reserved" in text_lower:
                if shape.top > (self.prs.slide_height * 0.8): return True
        return False
    
    def _is_decorative(self, shape) -> bool:
        try: return shape.element.nvSpPr.cNvPr.get("decorative") == "1"
        except AttributeError: return False

    def _determine_real_background(self, target_shp, slide) -> tuple:
        try:
            if target_shp.fill.type == MSO_FILL.SOLID and target_shp.fill.visible:
                alpha = target_shp.fill.fore_color.alpha
                if alpha is None or alpha == 1.0: return rgb_pptx_to_tuple(target_shp.fill.fore_color.rgb)
            elif target_shp.fill.type in (MSO_FILL.PICTURE, MSO_FILL.TEXTURE): return BG_IS_IMAGE
        except AttributeError: pass
        best_z = -1; best_bg_color = None; target_index = -1
        for i, s in enumerate(slide.shapes):
            if s.shape_id == target_shp.shape_id: target_index = i; break
        for i, candidate_shp in enumerate(slide.shapes):
            if i < target_index and shapes_overlap(target_shp, candidate_shp):
                if candidate_shp.shape_id == target_shp.shape_id: continue
                bg_color = None
                try:
                    if candidate_shp.shape_type == MSO_SHAPE_TYPE.TABLE: bg_color = BG_IS_TABLE 
                    elif candidate_shp.fill.type == MSO_FILL.SOLID:
                        if candidate_shp.fill.fore_color.alpha is None or candidate_shp.fill.fore_color.alpha == 1.0:
                            try: bg_color = rgb_pptx_to_tuple(candidate_shp.fill.fore_color.rgb)
                            except AttributeError: pass 
                    elif candidate_shp.fill.type in (MSO_FILL.PICTURE, MSO_FILL.TEXTURE): bg_color = BG_IS_IMAGE
                except AttributeError: pass
                if bg_color: best_bg_color = bg_color
        if best_bg_color: return best_bg_color
        return WHITE_RGB_TUPLE

    def _check_spelling(self, text, shape_name, slide_index, slide_issues):
        if not text or len(text) < 3: return
        text_no_url = re.sub(r'http\S+|www\.\S+|[\w\.-]+@[\w\.-]+', '', text)
        text_norm = text_no_url.replace("’", "'").replace("‘", "'"); text_spaced = re.sub(r'([a-z])([A-Z])', r'\1 \2', text_norm)
        text_spaced = re.sub(r'(\.)([A-Z])', r'\1 \2', text_spaced); text_spaced = re.sub(r'[-_/]', ' ', text_spaced)
        clean_text = re.sub(r"[^\w\s']", '', text_spaced); words = clean_text.split(); candidates = []
        for w in words:
            w_clean = w.strip("'")
            if not w_clean: continue
            if any(char.isdigit() for char in w_clean): continue
            if w_clean.isupper() and len(w_clean) > 1: continue
            if len(w_clean) > 3: candidates.append(w_clean)
        if not candidates: return
        unknown = self.spell.unknown(candidates)
        if unknown:
            safe_words = {"youtube", "linkedin", "tiktok", "instagram", "video", "intro", "outro", "agenda", "module"}
            true_typos = [w for w in unknown if w.lower() not in safe_words]
            if true_typos:
                typo_str = ", ".join(list(true_typos)[:5])
                slide_issues.append({"slide": slide_index, "check": "Copyediting - Spelling", "shape_name": shape_name, "result": "WARNING", "details": f"Potential typos found: {typo_str}..."})

    def _check_clarity_and_style(self, text, shape_name, slide_index, slide_issues, context="slide"):
        if not text or len(text.split()) < 4: return
        if self.clarity:
            res = self.clarity.check_reading_level(text, context)
            if res:
                res.update({"slide": slide_index, "check": "Clarity - Reading Level", "shape_name": shape_name})
                slide_issues.append(res)
            res = self.clarity.check_passive_voice(text, context)
            if res:
                res.update({"slide": slide_index, "check": "Clarity - Tone", "shape_name": shape_name})
                slide_issues.append(res)
            res = self.clarity.check_bad_habits(text)
            if res:
                res.update({"slide": slide_index, "check": "Brand Voice - Jargon", "shape_name": shape_name})
                slide_issues.append(res)

    def _check_reading_order(self, slide, slide_index, slide_issues):
        title_candidates = []
        if slide.shapes.title: title_candidates.append(slide.shapes.title)
        for shape in slide.shapes:
            if shape.is_placeholder and shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE):
                if shape not in title_candidates: title_candidates.append(shape)
        if not title_candidates:
            slide_issues.append({"slide": slide_index, "check": "Accessibility Reading Order", "shape_name": "Slide", "result": "FAIL", "details": "Slide missing standard Title placeholder."})
            return
        title_candidates.sort(key=lambda s: s.top); target_title = title_candidates[0] 
        for s in slide.shapes:
            if shapes_overlap(target_title, s) and not self._is_exempt_shape(s) and s.shape_id != target_title.shape_id:
                slide_issues.append({"slide": slide_index, "check": "Accessibility Reading Order", "shape_name": s.name, "result": "FAIL", "details": f"Object '{s.name}' obscures the Title."})

    def _check_hyperlinks(self, shape, slide_index, slide_issues):
        if not shape.has_text_frame: return
        for p in shape.text_frame.paragraphs:
            for r in p.runs:
                if r.hyperlink.address and not r.hyperlink.address.startswith("#"):
                    slide_issues.append({"slide": slide_index, "check": "Usability - Links", "shape_name": shape.name, "result": "INFO", "details": f"External Link found: {r.hyperlink.address}"})

    def _check_citations(self, text, shape_name, slide_index, slide_issues):
        facilitation_keywords = ["Talking Points:", "Instructional Activity:", "Do:", "Say:", "Transition:", "Instructor Guide", "Script:"]
        if any(k in text for k in facilitation_keywords): return 
        citation_pattern = r"\([A-Za-z\s&]+,\s?\d{4}\)"; source_header_pattern = r"(?i)(Source:|Sources:|Reference:|References:|Bibliography:|Works Cited:)"
        numeric_pattern = r"(\[\d+\]|\(\d+\)|Source\s+\d+)"
        if not (bool(re.search(citation_pattern, text)) or bool(re.search(source_header_pattern, text)) or bool(re.search(numeric_pattern, text))):
            if len(text.split()) > 50: 
                slide_issues.append({"slide": slide_index, "check": "Instructional Design - Citations", "shape_name": shape_name, "result": "WARNING", "details": "Large text block appears missing APA citation or Source header."})
    
    def _check_alt_text(self, shape, slide_index, slide_issues):
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and not self._is_decorative(shape):
            try: alt = shape.element.nvSpPr.cNvPr.get("descr", "")
            except: alt = ""
            if not alt: 
                # --- FIXED: Added Shape Name for Context ---
                slide_issues.append({"slide": slide_index, "check": "Accessibility - Alt Text", "shape_name": shape.name, "result": "FAIL", "details": f"Image '{shape.name}' is missing Alt Text."})

    def _check_fonts(self, shape, slide_index, slide_issues, is_ghost_title=False):
        if not shape.has_text_frame or not shape.text_frame.text.strip(): return
        if is_ghost_title: return
        is_title = shape.is_placeholder and shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip(): continue
                f_name = self._resolve_font_name(run.font.name or ("+mj-lt" if is_title else "+mn-lt"))
                f_size = run.font.size.pt if run.font.size else 0
                if is_title:
                    if f_name != TITLE_FONT_NAME: slide_issues.append({"slide": slide_index, "check": "Font Rules", "shape_name": shape.name, "result": "FAIL", "details": f"Title font is '{f_name}'."})
                else:
                    if f_size < BODY_FONT_SIZE_MIN and f_size > 0: slide_issues.append({"slide": slide_index, "check": "Font Rules", "shape_name": shape.name, "result": "FAIL", "details": f"Font size {f_size}pt below minimum."})

    def _check_brand_colors(self, shape, slide_index, slide_issues):
        try:
            if shape.fill.type == MSO_FILL.SOLID and shape.fill.visible:
                rgb = rgb_pptx_to_tuple(shape.fill.fore_color.rgb)
                if not is_brand_color(rgb): slide_issues.append({"slide": slide_index, "check": "Brand Consistency", "shape_name": shape.name, "result": "FAIL", "details": f"Non-brand object color found: {rgb}"})
        except: pass

    def _check_non_text_contrast(self, shape, slide, slide_index, slide_issues):
        if shape.has_text_frame and shape.text_frame.text.strip(): return
        try:
            if shape.fill.type == MSO_FILL.SOLID and shape.fill.visible:
                alpha = shape.fill.fore_color.alpha
                if alpha is not None and alpha < 1.0: return 
                obj_rgb = rgb_pptx_to_tuple(shape.fill.fore_color.rgb)
                bg_rgb = self._determine_real_background(shape, slide)
                if bg_rgb in (BG_IS_IMAGE, BG_IS_TABLE):
                    # Fixed in main loop to provide context
                    return
                ratio = calculate_contrast_ratio(obj_rgb, bg_rgb)
                if ratio < WCAG_MIN_GRAPHIC_RATIO:
                    slide_issues.append({"slide": slide_index, "check": "WCAG Graphic Contrast", "shape_name": shape.name, "result": "FAIL", "ratio": f"{ratio:.1f}:1", "details": f"Low contrast graphics. Ratio: {ratio:.1f}:1"})
        except: pass

    def _check_notes_formatting(self, slide, slide_index, slide_issues):
        if not slide.has_notes_slide: return
        notes_tf = slide.notes_slide.notes_text_frame
        if not notes_tf.text.strip(): return
        clean_text = self._strip_citations_for_analysis(notes_tf.text)
        self._check_single_quotes(clean_text, "Speaker Notes", slide_index, slide_issues)
        self._check_spelling(clean_text, "Speaker Notes", slide_index, slide_issues)
        self._check_citations(notes_tf.text, "Speaker Notes", slide_index, slide_issues)
        
        font_error_logged = False
        size_error_logged = False
        color_error_logged = False
        
        for paragraph in notes_tf.paragraphs:
            for run in paragraph.runs:
                if not run.text.strip(): continue
                
                # --- FIXED: Only check font name if user has configured one ---
                if not font_error_logged and NOTES_FONT_NAME and NOTES_FONT_NAME.strip():
                    if run.font.name is None: raw_font_name = "+mn-lt" 
                    else: raw_font_name = run.font.name
                    f_name = self._resolve_font_name(raw_font_name)
                    if f_name != NOTES_FONT_NAME: 
                        slide_issues.append({"slide": slide_index, "check": "Notes Formatting", "shape_name": "Speaker Notes", "result": "FAIL", "details": f"Notes font is '{raw_font_name}' (Should be {NOTES_FONT_NAME})"})
                        font_error_logged = True
                
                if not size_error_logged:
                    f_size = run.font.size.pt if run.font.size else 12
                    if not (NOTES_FONT_SIZE_MIN <= f_size <= NOTES_FONT_SIZE_MAX): 
                        slide_issues.append({"slide": slide_index, "check": "Notes Formatting", "shape_name": "Speaker Notes", "result": "FAIL", "details": f"Notes size is {f_size}pt (Should be {NOTES_FONT_SIZE_MIN}-{NOTES_FONT_SIZE_MAX})"})
                        size_error_logged = True
                if not color_error_logged:
                    try:
                        if run.font.color.rgb:
                            c_rgb = rgb_pptx_to_tuple(run.font.color.rgb)
                            if c_rgb != NOTES_FONT_COLOR_RGB: 
                                slide_issues.append({"slide": slide_index, "check": "Notes Formatting", "shape_name": "Speaker Notes", "result": "FAIL", "details": "Notes font color is not Black."})
                                color_error_logged = True
                    except AttributeError: pass

    def _check_notes_content(self, slide, slide_index, slide_issues):
        if not slide.has_notes_slide: return
        notes_tf = slide.notes_slide.notes_text_frame
        if not notes_tf.text.strip(): return
        notes_text = notes_tf.text; notes_lower = notes_text.lower()
        for header in REQUIRED_HEADERS:
            if header.lower() not in notes_lower:
                if header in ["Instructional Activity:", "Instructional Time:"]:
                    slide_issues.append({"slide": slide_index, "check": "Instructional Design", "shape_name": "Speaker Notes", "result": "FAIL", "details": f"Missing critical header: '{header}'"})

    def analyze_slide(self, slide, slide_index, is_exempt=False):
        self.total_slides_checked += 1
        content = self._extract_slide_text_content(slide)
        self.slide_content_map.append({
            "slide_number": slide_index, "title": content['title'], "full_text": content['text'], "notes": content['notes'], "visual_context": content['visual_context']
        })
        slide_issues = []
        if is_exempt:
            slide_issues.append({"slide": slide_index, "check": "Status", "shape_name": "Slide", "result": "EXEMPT", "details": "Slide excluded from analysis."})
            self.qa_report.extend(slide_issues)
            return

        self._check_reading_order(slide, slide_index, slide_issues)
        
        # --- FIXED: Use a set to prevent spamming the same error for one shape ---
        shapes_with_contrast_issues = set()
        
        for shape in slide.shapes:
            if self._is_exempt_shape(shape): continue
            is_ghost = False 
            self._check_fonts(shape, slide_index, slide_issues, is_ghost_title=is_ghost)
            self._check_brand_colors(shape, slide_index, slide_issues)
            self._check_non_text_contrast(shape, slide, slide_index, slide_issues)
            self._check_hyperlinks(shape, slide_index, slide_issues)
            self._check_alt_text(shape, slide_index, slide_issues)
            if shape.has_text_frame:
                if shape.text_frame.text.strip():
                    self._check_single_quotes(shape.text_frame.text, shape.name, slide_index, slide_issues)
                    self._check_spelling(shape.text_frame.text, shape.name, slide_index, slide_issues)
                    self._check_clarity_and_style(shape.text_frame.text, shape.name, slide_index, slide_issues)
                
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.text.strip():
                            try: text_rgb = rgb_pptx_to_tuple(run.font.color.rgb)
                            except: text_rgb = BLACK_RGB_TUPLE
                            bg_rgb = self._determine_real_background(shape, slide)
                            
                            # --- FIXED: Better Context for WCAG Text Over Image ---
                            if bg_rgb in (BG_IS_IMAGE, BG_IS_TABLE):
                                if shape.shape_id not in shapes_with_contrast_issues:
                                    snippet = run.text.strip().replace('\n', ' ')
                                    if len(snippet) > 50: snippet = snippet[:50] + "..."
                                    slide_issues.append({
                                        "slide": slide_index, 
                                        "check": "WCAG Text Contrast", 
                                        "shape_name": shape.name, 
                                        "result": "MANUAL REVIEW", 
                                        "details": f"Text '{snippet}' is over Image/Table. Verify contrast manually."
                                    })
                                    shapes_with_contrast_issues.add(shape.shape_id)
                                continue
                            
                            ratio = calculate_contrast_ratio(text_rgb, bg_rgb)
                            req = get_required_ratio(run)
                            if ratio < req:
                                txt_hex = rgb_to_hex(text_rgb)
                                bg_hex = rgb_to_hex(bg_rgb)
                                slide_issues.append({
                                    "slide": slide_index, 
                                    "check": "WCAG Text Contrast", 
                                    "shape_name": shape.name, 
                                    "result": "FAIL", 
                                    "ratio": f"{ratio:.1f}:1", 
                                    "details": f"Contrast {ratio:.1f}:1 fails WCAG {req}:1. Text: {txt_hex} on Bg: {bg_hex}."
                                })

        if slide.has_notes_slide:
            if slide.notes_slide.notes_text_frame.text.strip():
                self._check_notes_formatting(slide, slide_index, slide_issues)
                self._check_notes_content(slide, slide_index, slide_issues)
        
        self.qa_report.extend(slide_issues)

    def run_analysis(self):
        total_count = len(self.prs.slides)
        exempt_indices = set(EXEMPT_SPECIFIC_SLIDES)
        if EXEMPT_FIRST_SLIDE: exempt_indices.add(1)
        if EXEMPT_LAST_SLIDE: exempt_indices.add(total_count)
        
        for i, slide in enumerate(self.prs.slides):
            self.analyze_slide(slide, i + 1, is_exempt=(i+1 in exempt_indices))
        
        self._calculate_cadence_metrics()
        
        # Safe aggregation
        clean_report = [i for i in self.qa_report if isinstance(i, dict)]
        
        fails = set(i['slide'] for i in clean_report if i.get('result') == 'FAIL')
        score = 100.0
        if self.total_slides_checked > 0:
            score = round(((self.total_slides_checked - len(fails)) / self.total_slides_checked) * 100, 1)

        content_map = {str(s['slide_number']): s for s in self.slide_content_map}

        manual_reviews = len([i for i in clean_report if i.get('result') == 'MANUAL REVIEW'])
        
        reading_fails = len([i for i in clean_report if i.get('check') == 'Clarity - Reading Level'])
        passive_voice = len([i for i in clean_report if i.get('check') == 'Clarity - Tone'])
        jargon = len([i for i in clean_report if i.get('check') == 'Brand Voice - Jargon'])

        final_data = {
            "summary": {
                "presentation_name": self.filename,
                "date_generated": datetime.now().isoformat(),
                "total_slides_checked": self.total_slides_checked,
                "total_errors": len(clean_report),
                "manual_reviews": manual_reviews,
                "executive_metrics": {
                    "wcag_compliance_rate": score,
                    "slides_present_content": f"{self.gagne_metrics['Present Content']} ({round(self.gagne_time_distribution['Present Content'], 1)}%)",
                    "slides_elicit_performance": f"{self.gagne_metrics['Elicit Performance']} ({round(self.gagne_time_distribution['Elicit Performance'], 1)}%)"
                },
                "content_metrics": {
                    "reading_complexity_fails": reading_fails,
                    "passive_voice_count": passive_voice,
                    "jargon_count": jargon
                },
                "pacing_metrics": self.pacing_data
            },
            "detailed_issues": clean_report, 
            "slide_content": content_map
        }
        
        # --- HYBRID RETURN ---
        # Wrap dictionary in HybridReport so it behaves like a Dict for app.py
        # but like a List for qa_tool.py
        hybrid_response = HybridReport(final_data)
        return hybrid_response

    def export_debug_log(self, output_dir):
        csv_path = os.path.join(output_dir, f"Pacing_Debug_Log_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv")
        headers = ["Slide", "Type", "Activity?", "Sequence?", "Planned", "Implicit", "Projected", "Delta", "Funded?", "Notes"]
        try:
            with open(csv_path, 'w', newline='', encoding='utf-8') as f:
                writer = csv.DictWriter(f, fieldnames=headers)
                writer.writeheader()
                writer.writerows(self.debug_log)
            return csv_path
        except Exception as e:
            print(f"❌ Log Export Failed: {e}")
            return None
