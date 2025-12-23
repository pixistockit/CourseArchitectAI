import os
import copy
import re
import json
import logging
import math
from datetime import datetime
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.enum.dml import MSO_FILL
from pptx.util import Pt
from pptx.dml.color import RGBColor

# --- LOGGING SETUP ---
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger('fix_engine')

# --- ROBUST CONFIG IMPORT ---
try:
    from . import config as CFG
except ImportError:
    # Fallback if run directly or import fails
    class CFG:
        TITLE_FONT_NAME = "Calibri"
        BODY_FONT_NAME = "Calibri"
        NOTES_FONT_NAME = "Calibri"
        NOTES_FONT_COLOR_RGB = (0, 0, 0)
        WCAG_RATIO_NORMAL = 4.5
        BODY_FONT_SIZE_MIN = 18

def load_fixer_config():
    """Loads user-defined brand settings."""
    try:
        base_dir = os.path.dirname(os.path.abspath(__file__))
        config_path = os.path.join(base_dir, '../../data/config/brand_config.json')
        llm_config_path = os.path.join(base_dir, '../../data/config/llm_config.json')
        
        merged_config = {}
        if os.path.exists(config_path):
            with open(config_path, 'r') as f: merged_config.update(json.load(f))
        if os.path.exists(llm_config_path):
            with open(llm_config_path, 'r') as f: merged_config.update(json.load(f))
        return merged_config
    except Exception as e:
        logger.warning(f"Config load failed: {e}")
        return {}

class FixEngine:
    def __init__(self):
        self.user_config = load_fixer_config()
        
        # Resolve Settings (User > Config > Default)
        user_min = self.user_config.get('min_font_size')
        sys_min = getattr(CFG, 'BODY_FONT_SIZE_MIN', 18)
        self.min_font_size_pt = float(user_min) if user_min else float(sys_min)
        
        self.font_title = self.user_config.get('title_font') or getattr(CFG, 'TITLE_FONT_NAME', 'Calibri')
        self.font_body = self.user_config.get('body_font') or getattr(CFG, 'BODY_FONT_NAME', 'Calibri')
        self.font_notes = self.user_config.get('notes_font') or getattr(CFG, 'NOTES_FONT_NAME', 'Calibri')
        
        self.log_report = [] 

    def apply_fixes(self, source_pptx, fixes_payload, output_dir):
        logger.info(f"🚀 STARTING FIX ENGINE on {os.path.basename(source_pptx)}")
        try: 
            prs = Presentation(source_pptx)
        except Exception as e: 
            raise ValueError(f"Could not open PPTX: {e}")

        changes_made = 0
        self.log_report = []

        for fix in fixes_payload:
            try:
                slide_num = int(fix['slide_number'])
                target_type = fix.get('target', 'content')
                fix_type = fix.get('type', 'content_rewrite')
                target_shape_name = fix.get('shape_name')
                
                # Validation
                if slide_num < 1 or slide_num > len(prs.slides): 
                    continue
                
                slide = prs.slides[slide_num - 1]
                brand_rules = self._get_rules_for_slide(slide)

                # --- ROUTER: Apply Specific Fix Logic ---
                if fix_type == 'safe_compliance':
                    rule_name = fix.get('rule')
                    success = self._apply_safe_rule(slide, rule_name, brand_rules, slide_num, target_shape_name)
                    if success: changes_made += 1
                
                elif target_type in ['content', 'notes']:
                    new_text = fix.get('new_text', '')
                    target_text_frame = None
                    
                    if target_type == 'notes':
                        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                            target_text_frame = slide.notes_slide.notes_text_frame
                    else:
                        shape = self._find_best_text_shape(slide)
                        if shape: target_text_frame = shape.text_frame
                    
                    if target_text_frame:
                        layout_map = self._capture_layout_map(target_text_frame)
                        self._replace_text_enforced(target_text_frame, new_text, layout_map, target_type, brand_rules)
                        self.log_report.append(f"Slide {slide_num}: Rewrote {target_type} content.")
                        changes_made += 1
            
            except Exception as e:
                err = f"Error fixing Slide {slide_num}: {e}"
                logger.error(err)
                self.log_report.append(err)

        # Save Result
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        clean_name = os.path.splitext(os.path.basename(source_pptx))[0].replace(' ', '_')
        new_filename = f"{clean_name}_REMEDIATED_{timestamp}.pptx"
        output_path = os.path.join(output_dir, new_filename)
        
        prs.save(output_path)
        
        log_path = os.path.join(output_dir, f"fix_log_{timestamp}.json")
        with open(log_path, 'w') as f: 
            json.dump(self.log_report, f, indent=4)
            
        return output_path

    # --- HELPERS ---
    
    def _rgb_pptx_to_tuple(self, rgb):
        return (rgb[0], rgb[1], rgb[2])

    def _determine_real_background(self, target_shp, slide):
        # 1. Check if shape has own fill
        try:
            if target_shp.fill.type == MSO_FILL.SOLID and target_shp.fill.visible:
                alpha = target_shp.fill.fore_color.alpha
                if alpha is None or alpha == 1.0:
                    return self._rgb_pptx_to_tuple(target_shp.fill.fore_color.rgb)
        except: pass
        
        # Default to White if we can't determine (FixEngine is simpler than Analyzer)
        return (255, 255, 255) 

    def _iter_all_text_frames(self, shapes):
        for shape in shapes:
            name = shape.name
            if shape.has_text_frame: 
                yield shape.text_frame, name, shape 
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        if cell.text_frame: 
                            yield cell.text_frame, f"{name} (Table Cell)", shape
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from self._iter_all_text_frames(shape.shapes)

    def _apply_safe_rule(self, slide, rule_name, brand_rules, slide_num, target_shape_name=None):
        made_change = False
        
        def should_process(current_shape_name):
            if not target_shape_name: return True
            # Loose matching to handle minor naming differences
            return target_shape_name.strip() in current_shape_name.strip()

        if rule_name == 'fix_font_family':
            # Handle Speaker Notes specifically
            if target_shape_name == "Speaker Notes":
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    for p in slide.notes_slide.notes_text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.name != self.font_notes:
                                run.font.name = self.font_notes
                                made_change = True
                    if made_change: 
                        self.log_report.append(f"Slide {slide_num} (Notes): Standardized font to {self.font_notes}")
                return made_change
            
            # Handle Body/Title Shapes
            for text_frame, shape_name, shape_obj in self._iter_all_text_frames(slide.shapes):
                if not should_process(shape_name): continue
                
                target_font = self.font_body
                if shape_obj.is_placeholder:
                    ph_type = shape_obj.placeholder_format.type
                    if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE): 
                        target_font = self.font_title
                
                for p in text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.name != target_font:
                            run.font.name = target_font
                            made_change = True
                
                if made_change: 
                    self.log_report.append(f"Slide {slide_num} ({shape_name}): Standardized font to {target_font}")

        elif rule_name == 'fix_font_size':
            target_min = Pt(self.min_font_size_pt)
            for text_frame, shape_name, shape_obj in self._iter_all_text_frames(slide.shapes):
                if not should_process(shape_name): continue 
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size is not None and run.font.size < target_min:
                            old_sz = run.font.size.pt
                            run.font.size = target_min
                            self.log_report.append(f"Slide {slide_num} ({shape_name}): Bumped font {old_sz}pt -> {self.min_font_size_pt}pt")
                            made_change = True
                        elif len(run.text.strip()) > 0 and run.font.size is None:
                             run.font.size = target_min
                             self.log_report.append(f"Slide {slide_num} ({shape_name}): Enforced {self.min_font_size_pt}pt on unstyled text")
                             made_change = True
                                
        elif rule_name == 'fix_slide_title':
            if not slide.shapes.title: pass 
            elif not slide.shapes.title.text.strip():
                slide.shapes.title.text = "Insert Title Here"
                self.log_report.append(f"Slide {slide_num}: Inserted placeholder title.")
                made_change = True

        elif rule_name == 'fix_contrast':
            target_ratio = getattr(CFG, 'WCAG_RATIO_NORMAL', 4.5)
            for text_frame, shape_name, shape_obj in self._iter_all_text_frames(slide.shapes):
                if not should_process(shape_name): continue
                
                bg_rgb = self._determine_real_background(shape_obj, slide)
                
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        current_rgb = (0, 0, 0)
                        if run.font.color and run.font.color.type == MSO_SHAPE_TYPE.AUTO: 
                            current_rgb = (0,0,0)
                        elif run.font.color and run.font.color.rgb: 
                            current_rgb = (run.font.color.rgb[0], run.font.color.rgb[1], run.font.color.rgb[2])
                        
                        new_rgb = self._get_wcag_passing_color(current_rgb, bg_rgb, target_ratio)
                        if new_rgb != current_rgb:
                            run.font.color.rgb = RGBColor(new_rgb[0], new_rgb[1], new_rgb[2])
                            self.log_report.append(f"Slide {slide_num} ({shape_name}): Adjusted contrast {current_rgb} -> {new_rgb} (vs BG: {bg_rgb})")
                            made_change = True
        
        return made_change

    def _get_wcag_passing_color(self, fg_rgb, bg_rgb, target_ratio):
        def get_luminance(rgb):
            rgb_linear = []
            for c in rgb:
                c = c / 255.0
                val = c / 12.92 if c <= 0.03928 else ((c + 0.055) / 1.055) ** 2.4
                rgb_linear.append(val)
            return 0.2126 * rgb_linear[0] + 0.7152 * rgb_linear[1] + 0.0722 * rgb_linear[2]
        
        def get_contrast(rgb1, rgb2):
            lum1 = get_luminance(rgb1)
            lum2 = get_luminance(rgb2)
            return (max(lum1, lum2) + 0.05) / (min(lum1, lum2) + 0.05)
        
        current_fg = list(fg_rgb)
        if get_contrast(current_fg, bg_rgb) >= target_ratio: 
            return tuple(current_fg)
        
        bg_lum = get_luminance(bg_rgb)
        is_light_bg = bg_lum > 0.5
        
        for _ in range(20):
            contrast = get_contrast(current_fg, bg_rgb)
            if contrast >= target_ratio: 
                return tuple(current_fg)
            
            if is_light_bg:
                current_fg = [max(0, int(c * 0.90)) for c in current_fg] # Darken
            else:
                current_fg = [min(255, int(c + (255-c)*0.10)) for c in current_fg] # Lighten
                
        return (0,0,0) if is_light_bg else (255,255,255)

    def _get_rules_for_slide(self, slide):
        rules = {
            "body_font": getattr(CFG, 'BODY_FONT_NAME', 'Calibri'),
            "title_font": getattr(CFG, 'TITLE_FONT_NAME', 'Calibri'),
            "notes_font": getattr(CFG, 'NOTES_FONT_NAME', 'Calibri'),
            "safe_color": RGBColor(0, 0, 0)
        }
        
        # User Config Overrides
        if self.user_config.get('body_font'): rules['body_font'] = self.user_config['body_font']
        if self.user_config.get('title_font'): rules['title_font'] = self.user_config['title_font']
        if self.user_config.get('notes_font'): rules['notes_font'] = self.user_config['notes_font']
        
        return rules

    def _find_best_text_shape(self, slide):
        def iter_shapes(shapes_list):
            for s in shapes_list:
                if s.shape_type == MSO_SHAPE_TYPE.GROUP: yield from iter_shapes(s.shapes)
                else: yield s
        
        # Prioritize Body Placeholders
        for shape in slide.placeholders:
            try:
                if shape.placeholder_format.type in [PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT]:
                    if shape.has_text_frame: return shape
            except: pass
            
        # Fallback: Largest word count
        best_shape = None
        max_words = 0
        for shape in iter_shapes(slide.shapes):
            if shape.has_text_frame and shape.text.strip():
                word_count = len(shape.text.split())
                if word_count > max_words and word_count > 5: 
                    max_words = word_count
                    best_shape = shape
        return best_shape

    def _capture_layout_map(self, text_frame):
        layout_map = {}
        try:
            for p in text_frame.paragraphs:
                lvl = p.level
                if lvl not in layout_map:
                    layout_map[lvl] = {
                        "space_before": p.space_before, 
                        "space_after": p.space_after, 
                        "line_spacing": p.line_spacing, 
                        "alignment": p.alignment
                    }
        except: pass
        return layout_map

    def _replace_text_enforced(self, text_frame, new_content, layout_map, target_type, brand_rules):
        text_frame.clear()
        lines = new_content.split('\n')
        enforced_font = brand_rules["notes_font"] if target_type == 'notes' else brand_rules["body_font"]
        
        for line in lines:
            stripped = line.strip()
            if not stripped: continue
            
            p = text_frame.add_paragraph()
            level = 0
            clean_text = stripped
            
            # Simple list detection
            if line.startswith('    ') or line.startswith('\t'): level = 1
            if stripped.startswith(('o ', '- ', '– ')): 
                level = 1
                clean_text = re.sub(r'^[\s\t]*[o\-–]\s*', '', line).strip()
            elif stripped.startswith(('* ', '• ')): 
                level = 0
                clean_text = stripped[2:].strip()
            
            p.level = level
            p.text = clean_text
            
            # Restore Layout
            layout = layout_map.get(level, layout_map.get(0, {}))
            if layout:
                if layout["space_before"] is not None: p.space_before = layout["space_before"]
                if layout["space_after"] is not None: p.space_after = layout["space_after"]
                if layout["line_spacing"] is not None: p.line_spacing = layout["line_spacing"]
                if layout["alignment"] is not None: p.alignment = layout["alignment"]
            
            # Enforce Font
            if p.runs:
                run = p.runs[0]
                run.font.name = enforced_font 
                if not run.font.color or not run.font.color.rgb: 
                    run.font.color.rgb = brand_rules["safe_color"]
